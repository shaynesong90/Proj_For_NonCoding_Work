from pptx import Presentation
import os

def create_ppt(title, slides):
    prs = Presentation()

    for slide_data in slides:
        slide_layout = prs.slides.add_slide(prs.slide_layouts[1])
        slide_layout.shapes.title.text = slide_data["title"]
        slide_layout.placeholders[1].text = slide_data["content"]

    os.makedirs("outputs", exist_ok=True)
    output_path = os.path.join("outputs", f"{title}.pptx")
    prs.save(output_path)

    return output_path
